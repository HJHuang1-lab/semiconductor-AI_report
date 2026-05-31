import os
import sys
from datetime import datetime

# Set UTF-8 encoding for standard output
sys.stdout.reconfigure(encoding='utf-8')

print("正在初始化 Agentic Workflow 簡報生成器 (擴展版)...")

# ──────────────────────────────────────────────────────────────────────
# Step 1: Definition of Expanded Presentation Content (10 Content Slides)
# ──────────────────────────────────────────────────────────────────────
presentation_data = {
    "title": "從 Human SOP 到 Agentic Workflow",
    "subtitle": "基於任務拆解 (Task Decomposition) 的企業級 AI 系統構建指南",
    "presenter": f"AI 工作流研發小組 | {datetime.now().strftime('%Y-%m-%d')}",
    "slides": [
        {
            "title": "企業級 AI 應補齊的三個核心名詞",
            "subtitle": "Human SOP、Skill 與 Agentic Workflow 的本質區別與層級關係",
            "image_name": "agent_network.png",
            "layout_type": "split",  # Left cards, Right image
            "cards": [
                {
                    "title": "Human SOP (人類專用流程)",
                    "content": "傳統寫給人類看的文件。仰賴大腦「默會知識」自動補足脈絡與例外判定。對於 AI 來說是非結構化文字，理解與解析成本高，在複雜且多步驟執行中極易遺忘細節。"
                },
                {
                    "title": "Skill (Agent 的執行單元)",
                    "content": "將特定任務的方法論與踩坑紀錄打包。包含 Skill Markdown (行為指南)、References (範例與字典) 及 Scripts (傳統確定性處理)。專注於「單一任務」而非整條工作流。"
                },
                {
                    "title": "Agentic Workflow (生產線)",
                    "content": "將多個 Agents、Tools、Skills 與資料源串聯起來的完整工作流。如同工廠流水線，AI 夥伴各司其職，保證交付的高穩定性、高可觀測性與可修復性。"
                }
            ]
        },
        {
            "title": "為什麼 Mega Agent 必將失敗？",
            "subtitle": "企業生產環境 (Production-Ready) 的三大核心訴求",
            "image_name": "wafer_defect.png",
            "layout_type": "split",
            "cards": [
                {
                    "title": "買彩券式的隨機黑箱問題",
                    "content": "將龐大模糊任務整包丟給單一強大的 Mega Agent，其執行結果極具隨機性。一旦出錯，整坨輸入整坨輸出，中間推理步驟被深埋在 LLM 中，工程師根本無法進行局部的 Debug。"
                },
                {
                    "title": "生產級系統的三大支柱",
                    "content": "企業對於落地系統的硬性要求是穩定性 (Stability)、可觀測性 (Observability) 與可修復性 (Repairability)。看不見內部的 AI 黑箱是無法被允許上線服務的。"
                },
                {
                    "title": "分而治之 (Divide & Conquer)",
                    "content": "將任務拆解為多個小 Task，每個步驟定義極度明確的 Input/Output 與成功標準。各節點各司其職，出包時「壞哪裡改哪裡」，可局部迭代優化，降低系統複雜度。"
                }
            ]
        },
        {
            "title": "工作流轉換 - Step 1: 格式標準化",
            "subtitle": "將給人看的散文翻譯成 AI 能精確執行的合約與規範",
            "image_name": "yield_stack.png",
            "layout_type": "split",
            "cards": [
                {
                    "title": "參數化配置 (Parametrization)",
                    "content": "避免在 SOP 中寫死具體變數（如「使用 normal 模式」）。應將變數抽離為 parameters 參數（如 `mode`, `temperature` 等），使 SOP 轉化為高容錯、可重複調用的 Template 模板。"
                },
                {
                    "title": "引入 RFC 2119 規則強度規定",
                    "content": "用硬性規範規定 Agent 的行為：\n- MUST (硬性限制)：不可妥協，必須執行。\n- SHOULD (建議作法)：有強烈理由可例外，但須主動記錄 Log。\n- MAY (彈性選擇)：Agent 可基於 context 自主決定。"
                },
                {
                    "title": "結構化 Markdown 區塊切分",
                    "content": "使用 Markdown 清晰劃分 Parameters、Steps 與 Error Handling 等區區塊。不僅便於人類閱讀，更能直接塞入 MCP 等標準接口當作 Agent 的行為特徵規範。"
                }
            ]
        },
        {
            "title": "工作流轉換 - Step 2: 任務拆解與連結",
            "subtitle": "精確劃分 Skill 的邊界，並使用 Artifacts 建立透明管道",
            "image_name": "metrology_sensor.png",
            "layout_type": "split",
            "cards": [
                {
                    "title": "劃定 Skill 防守邊界 (Boundary)",
                    "content": "技能的範圍大小是關鍵。防守範圍太大，Agent 容易樣樣通但做得差強人意；防守範圍太小，會變成每走一步都要聽人類指示。應以「具備獨立產出」作為單一 Skill 的邊界。"
                },
                {
                    "title": "建立獨立的管道節點 (Pipeline Steps)",
                    "content": "每個步驟是一個獨立的 Skill 或獨立的 Agent，有其專屬的 SOP。若某步驟有 bug，只需修改該步驟的 SOP，後續的機器設定或判定邏輯完全不受影響，大幅降低複合誤差。"
                },
                {
                    "title": "以 Artifacts (JSON) 實體數據串接",
                    "content": "節點與節點之間不使用模糊的口頭語言或「心電感應」，而是靠清楚定義的 Artifacts（如 JSON 檔案）傳遞。前一個節點的 Output 完美對接下一個節點的 Input。"
                }
            ]
        },
        {
            "title": "工作流轉換 - Step 3: 雙向開發與迭代",
            "subtitle": "用敏捷 Scrum 精神快速收斂並捕捉大腦中的默會知識",
            "image_name": "yield_graph.png",
            "layout_type": "split",
            "cards": [
                {
                    "title": "暴露「默會知識 (Tacit Knowledge)」",
                    "content": "指存在於人類大腦記憶中、難以言表且自己不易察覺的潛規則（如「羊毛衣不可烘乾」）。寫第一版 SOP 時往往會自動忽略默會知識，必須讓 Agent 跑出去撞牆才能暴露出這些漏洞。"
                },
                {
                    "title": "小步快跑的敏捷迭代循環",
                    "content": "世界上不存在閉門造車就能完美的 SOP。應秉持敏捷精神：花兩天產出粗糙 SOP -> 實跑 50 次暴露出錯點 -> 回頭在 SOP 補上 MUST/SHOULD 規則 -> 再跑，迅速收斂系統漏洞。"
                },
                {
                    "title": "迭代效率大於理論完美",
                    "content": "與其花兩個月構思想像中的完美情境，不如兩天快速上線，並在兩週內跑五十次實例 iteration。真實執行反饋是排除 Edge Cases、建立生產級系統的唯一法寶。"
                }
            ]
        },
        {
            "title": "工作流轉換 - Step 4: 整合與執行環境",
            "subtitle": "利用 MCP 開放協定統一接口，並設立人工安全防線",
            "image_name": "litho_system.png",
            "layout_type": "split",
            "cards": [
                {
                    "title": "MCP 協定：AI 世界的 USB-C",
                    "content": "Model Context Protocol (MCP) 是由 Anthropic 捐贈給 Linux 基金會的開放協定。它統一了 LLM/Agent 與外部 tools、resources 的調用規格，免去針對不同平台重複開發整合的痛苦。"
                },
                {
                    "title": "人工確認點 (Human-in-the-Loop)",
                    "content": "任何成熟的 Agent 都有無法處理的 edge cases。必須在高風險決策前設計 Checkpoint（如大額撥款、主權限變更），強迫 Agent 暫停並等待人類核准，確保風險 100% 可控。"
                },
                {
                    "title": "人類掌舵，AI 划槳",
                    "content": "整合的終極目的不是建立失控的自主機器人，而是建立一個人類擁有最終裁決權、但所有機械、確定性與重複性的繁瑣工作均由 AI 高效代勞的完美閉環系統。"
                }
            ]
        },
        {
            "title": "Triage 案例 - 1. 標準化與任務拆解",
            "subtitle": "針對 200 人公司日常非結構化雜事請求 triage 流程的拆解實踐",
            "image_name": "chamber_control.png",
            "layout_type": "split",
            "cards": [
                {
                    "title": "INTERNAL_REQUEST_TRIAGE SOP",
                    "content": "參數化輸入：`source_channel`, `raw_text`, `employee_id`。\n規則約束：\n- MUST 驗證員工在職狀態。\n- MUST 將請求分類成 IT、HR、Finance 三大類。\n- SHOULD 基於關鍵字判斷優先級。"
                },
                {
                    "title": "拆解為兩個獨立技能",
                    "content": "將流程攔腰折半，拆成兩個職責單一的技能：\n1. `request-triage`：專攻員工身分驗證、分類、assignee 推薦與優先級判定，產出 JSON Artifact。\n2. `reply-drafting`：專攻讀取 JSON 資訊並草擬親切回覆。"
                },
                {
                    "title": "澄清機制與 JSON 中介檔串接",
                    "content": "若 triage 節點發現內容過於模糊，會在 JSON 中設置 `needs_clarification=true`，技能 B 讀取後將 `MUST` 額外草擬 2-3 個澄清問題。兩個 Skill 各自為戰，互不干擾，方便維護與調整口吻。"
                }
            ]
        },
        {
            "title": "Triage 案例 - 2. 迭代調優與真實整合",
            "subtitle": "從 edge case 踩坑中完善防守，並對接企業真實生產系統",
            "image_name": "fab_agv.png",
            "layout_type": "split",
            "cards": [
                {
                    "title": "雙向迭代與邊界收斂",
                    "content": "實跑初期，發現很多財務報銷 ticket 被誤分到 Other，或者 assignee 老是推薦給已離職員工。我們回頭補上最新員工庫匹配規則與財務關鍵字對照表，幾輪迭代後 triage 準確度即達 98%。"
                },
                {
                    "title": "對接真實 MCP 工具伺服器",
                    "content": "透過 Google Sheet MCP 與 Slack MCP。自動將 request 分類結果、assignee、優先級與處理草稿登錄到內部的 tracking sheet，並主動發送 Slack 動態通知給具體負責人。"
                },
                {
                    "title": "設立安全 Checkpoint 防線",
                    "content": "設計 Human-in-the-Loop checkpoint。當 category 為 Finance 且涉及報銷金額大於 $5000，或涉及 IT 主管理員權限變更時，強制暫停並彈出 Slack 確認視窗，待人類點擊 Approve 後才寫入系統。"
                }
            ]
        },
        {
            "title": "未來 AI 工作者的核心競爭力",
            "subtitle": "學會「如何使用 AI」半年就過時，學會「為 AI 設計工作流」越來越值錢",
            "image_name": "cleanroom_engineer.png",
            "layout_type": "standard",  # Three beautiful vertical cards
            "cards": [
                {
                    "title": "挑選你最不想做的那個 SOP",
                    "content": "不用一次想把整間公司流程都自動化，那會把自己搞死。從自己手邊最無聊、最討厭且不斷重複的 Human SOP（週報撰寫、新人 onboard、上線前 checklist）開始動手，先做出能省去 30% 時間的版本。"
                },
                {
                    "title": "從「AI 用戶」晉升為「AI 架構師」",
                    "content": "未來 MCP 與 Multi-Agent 會像空氣一樣普及。只會對著對話框提問 (Prompting) 的能力貶值極快。能看透業務本質、將複雜流程拆解並設計成 Agent 可穩定執行的 workflow 專家，其身價將只增不減。"
                },
                {
                    "title": "敏捷迭代大於理論完美",
                    "content": "在企業級 Agentic Workflow 實踐中，花費兩個月閉門造車產出的「完美SOP」往往一上線就崩潰。唯有敏捷開發、小步快跑、一週跑五十次 iteration 快速收斂默會知識，才是讓 AI 穩定落地生產環境的唯一法寶。"
                }
            ]
        }
    ]
}

# ──────────────────────────────────────────────────────────────────────
# Step 2: Widescreen Presentation Generation using python-pptx
# ──────────────────────────────────────────────────────────────────────
print("正在生成高質感、深色科技風簡報 (11頁完整版)...")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("❌ 錯誤：未安裝 python-pptx 套件，請執行 pip install python-pptx")
    sys.exit(1)

prs = Presentation()
# Set aspect ratio to 16:9 widescreen (13.333 x 7.5 Inches)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Premium Color Palette Definitions (Consistent with yield presentation style)
COLOR_BG = RGBColor(15, 23, 42)          # Deep Slate-Black
COLOR_CARD_BG = RGBColor(30, 41, 59)     # Lighter Slate Blue for cards
COLOR_CARD_BORDER = RGBColor(51, 65, 85)  # Thin card borders
COLOR_TITLE = RGBColor(248, 250, 252)    # Bright White
COLOR_SUBTITLE = RGBColor(56, 189, 248)  # Neon Cyan
COLOR_GOLD = RGBColor(253, 224, 71)      # Accent Gold for highlights
COLOR_TEXT = RGBColor(203, 213, 225)      # Soft Light Gray
COLOR_MUTED = RGBColor(148, 163, 184)     # Slate Gray for metadata

assets_dir = "assets"

def apply_background(slide, color):
    """Draws a full-slide rectangle to apply solid background color without standard borders."""
    bg = slide.shapes.add_shape(
        1,  # 1 represents MSO_SHAPE.RECTANGLE
        0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.color.rgb = color
    bg.line.width = Pt(0)

# ── 1. Create Title Slide ─────────────────────────────────────────────
slide_layout = prs.slide_layouts[6] # Blank layout
title_slide = prs.slides.add_slide(slide_layout)
apply_background(title_slide, COLOR_BG)

# Use semiconductor_ai.png as title slide cover image
img_cover_default = os.path.join(assets_dir, "semiconductor_ai.png")

if os.path.exists(img_cover_default):
    # Left side: Text card with titles
    accent_card = title_slide.shapes.add_shape(
        1, Inches(0.8), Inches(1.5), Inches(6.0), Inches(4.5)
    )
    accent_card.fill.solid()
    accent_card.fill.fore_color.rgb = COLOR_CARD_BG
    accent_card.line.color.rgb = COLOR_CARD_BORDER
    accent_card.line.width = Pt(1.5)
    
    title_box = title_slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(5.2), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p1 = tf.paragraphs[0]
    p1.text = presentation_data["title"]
    p1.font.name = "Microsoft JhengHei"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TITLE
    p1.space_after = Pt(14)
    
    p2 = tf.add_paragraph()
    p2.text = presentation_data["subtitle"]
    p2.font.name = "Microsoft JhengHei"
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_SUBTITLE
    
    meta_box = title_slide.shapes.add_textbox(Inches(1.2), Inches(4.8), Inches(5.2), Inches(1.0))
    tf_meta = meta_box.text_frame
    tf_meta.word_wrap = True
    tf_meta.margin_left = tf_meta.margin_right = tf_meta.margin_top = tf_meta.margin_bottom = 0
    
    p_meta = tf_meta.paragraphs[0]
    p_meta.text = presentation_data["presenter"]
    p_meta.font.name = "Segoe UI"
    p_meta.font.size = Pt(12)
    p_meta.font.color.rgb = COLOR_MUTED
    
    # Right side: Visual card with the glowing cover image
    img_card = title_slide.shapes.add_shape(
        1, Inches(7.2), Inches(1.5), Inches(5.333), Inches(4.5)
    )
    img_card.fill.solid()
    img_card.fill.fore_color.rgb = COLOR_CARD_BG
    img_card.line.color.rgb = COLOR_CARD_BORDER
    img_card.line.width = Pt(1.5)
    
    img_pad = Inches(0.15)
    title_slide.shapes.add_picture(
        img_cover_default, 
        Inches(7.2) + img_pad, Inches(1.5) + img_pad, 
        Inches(5.333) - (img_pad * 2), Inches(4.5) - (img_pad * 2)
    )
else:
    # Full screen text card layout if image is missing
    accent_card = title_slide.shapes.add_shape(
        1, Inches(1.0), Inches(1.5), Inches(11.333), Inches(4.5)
    )
    accent_card.fill.solid()
    accent_card.fill.fore_color.rgb = COLOR_CARD_BG
    accent_card.line.color.rgb = COLOR_CARD_BORDER
    accent_card.line.width = Pt(1.5)
    
    title_box = title_slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.333), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = presentation_data["title"]
    p1.font.name = "Microsoft JhengHei"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TITLE
    p1.space_after = Pt(24)
    
    p2 = tf.add_paragraph()
    p2.text = presentation_data["subtitle"]
    p2.font.name = "Microsoft JhengHei"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_SUBTITLE
    
    meta_box = title_slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10.333), Inches(1.0))
    tf_meta = meta_box.text_frame
    p_meta = tf_meta.paragraphs[0]
    p_meta.text = presentation_data["presenter"]
    p_meta.font.name = "Segoe UI"
    p_meta.font.size = Pt(13)
    p_meta.font.color.rgb = COLOR_MUTED

# ── 2. Create Content Slides ──────────────────────────────────────────
def add_split_slide(s_data, img_path):
    """Generates a NotebookLM-style split slide: Left 3 cards, Right 1 beautiful AI Image."""
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide, COLOR_BG)
    
    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.333), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p_title = tf.paragraphs[0]
    p_title.text = s_data["title"]
    p_title.font.name = "Microsoft JhengHei"
    p_title.font.size = Pt(26)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_SUBTITLE
    p_title.space_after = Pt(6)
    
    p_sub = tf.add_paragraph()
    p_sub.text = s_data["subtitle"]
    p_sub.font.name = "Microsoft JhengHei"
    p_sub.font.size = Pt(13)
    p_sub.font.color.rgb = COLOR_MUTED
    
    # Left Column: Structured Cards (NotebookLM deep analysis style)
    cards = s_data["cards"]
    card_x = Inches(1.0)
    card_width = Inches(5.3)
    
    # Adjust card heights based on count (usually 3 cards in the expanded version)
    num_cards = len(cards)
    if num_cards == 1:
        card_positions = [(Inches(1.8), Inches(4.8), cards[0])]
    elif num_cards == 2:
        card_h = Inches(2.25)
        gap = Inches(0.3)
        card_positions = [
            (Inches(1.8), card_h, cards[0]),
            (Inches(1.8) + card_h + gap, card_h, cards[1])
        ]
    else: # 3 cards
        card_h = Inches(1.45)
        gap = Inches(0.18)
        card_positions = [
            (Inches(1.8), card_h, cards[0]),
            (Inches(1.8) + card_h + gap, card_h, cards[1]),
            (Inches(1.8) + (card_h + gap) * 2, card_h, cards[2])
        ]
        
    for card_y, card_h, card in card_positions:
        card_shape = slide.shapes.add_shape(1, card_x, card_y, card_width, card_h)
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = COLOR_CARD_BG
        card_shape.line.color.rgb = COLOR_CARD_BORDER
        card_shape.line.width = Pt(1.5)
        
        pad = Inches(0.18)
        text_box = slide.shapes.add_textbox(card_x + pad, card_y + pad, card_width - (pad * 2), card_h - (pad * 2))
        tf_card = text_box.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = tf_card.margin_right = tf_card.margin_top = tf_card.margin_bottom = 0
        
        p_c_title = tf_card.paragraphs[0]
        p_c_title.text = f"📍 {card['title']}"
        p_c_title.font.name = "Microsoft JhengHei"
        p_c_title.font.size = Pt(12 if num_cards == 3 else 14)
        p_c_title.font.bold = True
        p_c_title.font.color.rgb = COLOR_GOLD
        p_c_title.space_after = Pt(4 if num_cards == 3 else 6)
        
        p_c_body = tf_card.add_paragraph()
        p_c_body.text = card["content"]
        p_c_body.font.name = "Microsoft JhengHei"
        p_c_body.font.size = Pt(9.5 if num_cards == 3 else 11)
        p_c_body.font.color.rgb = COLOR_TEXT
        p_c_body.line_spacing = 1.2
        
    # Right Column: Visual Image Card
    img_x = Inches(6.8)
    img_y = Inches(1.8)
    img_width = Inches(5.5)
    img_height = Inches(4.8)
    
    img_card = slide.shapes.add_shape(1, img_x, img_y, img_width, img_height)
    img_card.fill.solid()
    img_card.fill.fore_color.rgb = COLOR_CARD_BG
    img_card.line.color.rgb = COLOR_CARD_BORDER
    img_card.line.width = Pt(1.5)
    
    img_pad = Inches(0.15)
    if os.path.exists(img_path):
        slide.shapes.add_picture(
            img_path, 
            img_x + img_pad, img_y + img_pad, 
            img_width - (img_pad * 2), img_height - (img_pad * 2)
        )
    else:
        text_box = slide.shapes.add_textbox(img_x + Inches(0.5), img_y + Inches(2.0), img_width - Inches(1.0), Inches(1.0))
        tf_img = text_box.text_frame
        p_img = tf_img.paragraphs[0]
        p_img.text = f"[ 📊 晶圓製造製程示意圖 ]\n({os.path.basename(img_path)})"
        p_img.font.name = "Microsoft JhengHei"
        p_img.font.size = Pt(14)
        p_img.font.color.rgb = COLOR_MUTED
        p_img.alignment = 1

def add_standard_slide(s_data):
    """Generates a standard multi-card widescreen slide for layout diversity."""
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide, COLOR_BG)
    
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.333), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p_title = tf.paragraphs[0]
    p_title.text = s_data["title"]
    p_title.font.name = "Microsoft JhengHei"
    p_title.font.size = Pt(26)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_SUBTITLE
    p_title.space_after = Pt(6)
    
    p_sub = tf.add_paragraph()
    p_sub.text = s_data["subtitle"]
    p_sub.font.name = "Microsoft JhengHei"
    p_sub.font.size = Pt(13)
    p_sub.font.color.rgb = COLOR_MUTED
    
    cards = s_data["cards"]
    num_cards = len(cards)
    
    content_top = Inches(1.8)
    content_height = Inches(4.8)
    
    if num_cards == 2:
        card_width = Inches(5.4)
        card_gap = Inches(0.533)
    elif num_cards == 3:
        card_width = Inches(3.5)
        card_gap = Inches(0.416)
    else:
        card_width = Inches(11.333)
        card_gap = Inches(0)
        
    for idx, card in enumerate(cards):
        left_pos = Inches(1.0) + idx * (card_width + card_gap)
        
        card_shape = slide.shapes.add_shape(1, left_pos, content_top, card_width, content_height)
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = COLOR_CARD_BG
        card_shape.line.color.rgb = COLOR_CARD_BORDER
        card_shape.line.width = Pt(1.5)
        
        pad = Inches(0.3)
        text_box = slide.shapes.add_textbox(left_pos + pad, content_top + pad, card_width - (pad * 2), content_height - (pad * 2))
        tf_card = text_box.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = tf_card.margin_right = tf_card.margin_top = tf_card.margin_bottom = 0
        
        p_c_title = tf_card.paragraphs[0]
        p_c_title.text = f"📍 {card['title']}"
        p_c_title.font.name = "Microsoft JhengHei"
        p_c_title.font.size = Pt(15)
        p_c_title.font.bold = True
        p_c_title.font.color.rgb = COLOR_GOLD
        p_c_title.space_after = Pt(12)
        
        p_c_body = tf_card.add_paragraph()
        p_c_body.text = card["content"]
        p_c_body.font.name = "Microsoft JhengHei"
        p_c_body.font.size = Pt(11)
        p_c_body.font.color.rgb = COLOR_TEXT
        p_c_body.line_spacing = 1.3
        
# Build Slides dynamically
for slide_data in presentation_data["slides"]:
    layout_type = slide_data.get("layout_type", "split")
    img_name = slide_data.get("image_name", "")
    img_path = os.path.join(assets_dir, img_name)
    
    if layout_type == "split":
        print(f"正在生成 Split 投影片: 『{slide_data['title']}』 (圖片: {img_name})")
        add_split_slide(slide_data, img_path)
    else:
        print(f"正在生成 Standard 投影片: 『{slide_data['title']}』")
        add_standard_slide(slide_data)

# ── 3. Create Thank You / End Slide ───────────────────────────────────
end_slide = prs.slides.add_slide(slide_layout)
apply_background(end_slide, COLOR_BG)

end_card = end_slide.shapes.add_shape(
    1, Inches(2.0), Inches(2.0), Inches(9.333), Inches(3.5)
)
end_card.fill.solid()
end_card.fill.fore_color.rgb = COLOR_CARD_BG
end_card.line.color.rgb = COLOR_CARD_BORDER
end_card.line.width = Pt(1.5)

end_box = end_slide.shapes.add_textbox(Inches(2.5), Inches(2.6), Inches(8.333), Inches(2.3))
tf_end = end_box.text_frame
tf_end.word_wrap = True
tf_end.margin_left = tf_end.margin_right = tf_end.margin_top = tf_end.margin_bottom = 0

p_end1 = tf_end.paragraphs[0]
p_end1.text = "簡報結束，謝謝觀看"
p_end1.font.name = "Microsoft JhengHei"
p_end1.font.size = Pt(36)
p_end1.font.bold = True
p_end1.font.color.rgb = COLOR_TITLE
p_end1.alignment = 1 # Centered
p_end1.space_after = Pt(14)

p_end2 = tf_end.add_paragraph()
p_end2.text = "From Human SOP to Agentic Workflow Presentation"
p_end2.font.name = "Segoe UI"
p_end2.font.size = Pt(14)
p_end2.font.color.rgb = COLOR_SUBTITLE
p_end2.alignment = 1 # Centered

# Save PPTX with high fault tolerance for file locks (PermissionError)
output_ppt = "SOP_to_Agentic_Workflow_Presentation.pptx"
try:
    prs.save(output_ppt)
    print(f"🎉 成功生成並儲存 PPT 檔案於: {output_ppt}")
except PermissionError:
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    backup_ppt = f"SOP_to_Agentic_Workflow_Presentation_{timestamp}.pptx"
    print(f"⚠️ 警告：無法寫入 {output_ppt}。這通常是因為該簡報檔目前已在 PowerPoint 中開啟。")
    print(f"👉 啟動安全備份儲存，檔案已儲存為: {backup_ppt}")
    prs.save(backup_ppt)
    print("💡 提示：若要更新主簡報檔，請先關閉您電腦上的 PowerPoint 程式，然後重新執行本腳本。")
